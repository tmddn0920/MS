from lxml import etree
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Emu, Inches, Pt
from pptx.oxml.ns import nsmap


FONT = "Malgun Gothic"
BLACK = RGBColor(0x14, 0x14, 0x14)
GOLD = RGBColor(0xB8, 0x96, 0x5A)
INK = RGBColor(0x2A, 0x27, 0x24)
MUTED = RGBColor(0x6E, 0x68, 0x61)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
SOFT = RGBColor(0xF6, 0xF5, 0xF3)
LINE = RGBColor(0xE4, 0xE0, 0xDA)
DARK = RGBColor(0x1C, 0x1B, 0x19)


def set_run(run, size, bold=False, color=INK, font=FONT, italic=False):
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    rPr = run._r.get_or_add_rPr()
    for tag in ("ea", "cs"):
        el = rPr.find(qn(f"a:{tag}"))
        if el is None:
            el = etree.SubElement(rPr, qn(f"a:{tag}"))
        el.set("typeface", font)


def add_text(slide, l, t, w, h, text, size=16, bold=False, color=INK, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    tf.paragraphs[0].alignment = align
    try:
        box.text_frame._txBody.bodyPr.set("anchor", {MSO_ANCHOR.TOP: "t", MSO_ANCHOR.MIDDLE: "ctr", MSO_ANCHOR.BOTTOM: "b"}[anchor])
    except Exception:
        pass
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    set_run(run, size, bold=bold, color=color)
    return box


def add_lines(slide, l, t, w, h, lines, size=15, color=INK, spacing=8):
    box = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(spacing)
        if isinstance(item, tuple):
            text, is_bold, c, s = item
        else:
            text, is_bold, c, s = item, False, color, size
        run = p.add_run()
        run.text = text
        set_run(run, s, bold=is_bold, color=c)
    return box


def fill(shape, color):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def fill_line(shape, color, line_color=None):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    if line_color is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(1)


def rect(slide, l, t, w, h, color, line_color=None):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(h))
    fill_line(s, color, line_color)
    s.shadow.inherit = False
    return s


def round_rect(slide, l, t, w, h, color, line_color=None):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(h))
    fill_line(s, color, line_color)
    s.adjustments[0] = 0.08
    s.shadow.inherit = False
    return s


def labeled_box(slide, l, t, w, h, title, body, fill_color=SOFT, title_color=BLACK, body_color=MUTED):
    shape = round_rect(slide, l, t, w, h, fill_color)
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.16)
    tf.margin_right = Inches(0.16)
    tf.margin_top = Inches(0.16)
    tf.margin_bottom = Inches(0.12)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    r = p.add_run()
    r.text = title
    set_run(r, 13, bold=True, color=title_color)
    p2 = tf.add_paragraph()
    p2.space_before = Pt(6)
    r2 = p2.add_run()
    r2.text = body
    set_run(r2, 12, color=body_color)
    return shape


def footer(slide, page, total=8):
    rect(slide, 0, 7.32, 13.333, 0.18, DARK)
    add_text(slide, 0.5, 7.33, 8, 0.16, "MCM AI Concierge  ·  RAG Pipeline", 9, color=GOLD, anchor=MSO_ANCHOR.MIDDLE)
    add_text(slide, 11.2, 7.33, 1.6, 0.16, f"{page}  /  {total}", 9, color=WHITE, align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)


def header(slide, kicker, title):
    rect(slide, 0, 0, 0.12, 7.5, GOLD)
    add_text(slide, 0.55, 0.28, 12, 0.28, kicker, 12, bold=True, color=GOLD)
    add_text(slide, 0.55, 0.52, 12, 0.5, title, 28, bold=True, color=BLACK)
    rect(slide, 0.55, 1.08, 1.4, 0.04, GOLD)


def new_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    return slide


def add_arrow(slide, l, t):
    add_text(slide, l, t, 0.35, 0.4, "→", 22, bold=True, color=GOLD, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


def main():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # --------------------------------------------------
    # 1. Title
    # --------------------------------------------------
    s = new_slide(prs)
    rect(s, 0, 0, 13.333, 7.5, DARK)
    rect(s, 0, 0, 0.16, 7.5, GOLD)
    add_text(s, 0.8, 1.7, 11, 0.3, "PRODUCT RAG  ·  AI AGENT", 14, bold=True, color=GOLD)
    add_text(s, 0.8, 2.1, 11.5, 1.1, "MCM AI Concierge", 48, bold=True, color=WHITE)
    add_text(
        s, 0.8, 3.3, 10.5, 0.7,
        "제품 Database를 RAG 문서로 만들고, 벡터 검색으로 관련 정보를 찾아\nAI Agent가 답변하는 구조입니다.",
        18, color=RGBColor(0xD6, 0xD0, 0xC6),
    )

    steps = [
        ("01", "Database"),
        ("02", "RAG"),
        ("03", "Vector"),
        ("04", "Search"),
        ("05", "Agent"),
    ]
    x = 0.8
    for i, (num, name) in enumerate(steps):
        add_text(s, x, 5.35, 1.6, 0.22, num, 11, bold=True, color=GOLD)
        add_text(s, x, 5.58, 1.6, 0.32, name, 18, bold=True, color=WHITE)
        if i < len(steps) - 1:
            add_text(s, x + 1.55, 5.5, 0.4, 0.35, "→", 16, color=GOLD)
        x += 2.05
    add_text(s, 0.8, 6.85, 10, 0.25, "현재 OpenAI  ·  이후 Qwen 확장 가능", 12, color=RGBColor(0xA9, 0xA3, 0x9A))

    # --------------------------------------------------
    # 2. Folder structure
    # --------------------------------------------------
    s = new_slide(prs)
    header(s, "STRUCTURE", "폴더 구조")
    footer(s, 2)

    folders = [
        ("database/", "제품 원본 데이터", "mcm_products_v2.json"),
        ("rag/", "RAG 문서", "documents.jsonl"),
        ("vector_db/", "벡터 저장소", "documents.json  ·  index.faiss"),
        ("config/", "LLM 설정", "openai_config.json"),
        ("qwen/", "Qwen 확장", "chat.py  ·  chat_ui.py"),
    ]
    x = 0.55
    for name, role, files in folders:
        round_rect(s, x, 1.4, 2.35, 2.15, SOFT)
        add_text(s, x + 0.14, 1.55, 2.1, 0.45, name, 16, bold=True, color=BLACK)
        add_text(s, x + 0.14, 2.05, 2.1, 0.5, role, 13, color=GOLD)
        add_text(s, x + 0.14, 2.55, 2.1, 0.75, files, 11, color=MUTED)
        x += 2.5

    scripts = [
        ("build_rag_documents.py", "제품을 검색용 RAG 문서로 변환"),
        ("build_vector_db.py", "문서를 embedding 후 FAISS 벡터로 저장"),
        ("chat_openai.py", "검색 + OpenAI Agent 답변"),
        ("api.py", "FastAPI 스트리밍 엔드포인트"),
        ("chat_ui_openai.py", "OpenAI 대화 UI"),
        ("test_search.py", "검색 품질 확인"),
    ]
    y = 3.8
    left = scripts[:3]
    right = scripts[3:]
    for i, (name, desc) in enumerate(left):
        add_text(s, 0.55, y + i * 0.85, 5.8, 0.28, name, 14, bold=True, color=BLACK)
        add_text(s, 0.55, y + 0.3 + i * 0.85, 5.8, 0.28, desc, 13, color=MUTED)
    for i, (name, desc) in enumerate(right):
        add_text(s, 7.1, y + i * 0.85, 5.6, 0.28, name, 14, bold=True, color=BLACK)
        add_text(s, 7.1, y + 0.3 + i * 0.85, 5.6, 0.28, desc, 13, color=MUTED)

    # --------------------------------------------------
    # 3. Pipeline
    # --------------------------------------------------
    s = new_slide(prs)
    header(s, "FLOW", "전체 파이프라인")
    footer(s, 3)

    pipeline = [
        ("1", "Product DB", "공식 제품 정보를\nJSON으로 보관"),
        ("2", "RAG 문서", "주제별 문서로\n잘게 나눔"),
        ("3", "벡터화", "의미를 숫자 벡터로\n변환해 저장"),
        ("4", "검색", "질문과 벡터를\n비교해 관련 문서 선택"),
        ("5", "AI Agent", "찾은 문서를 근거로\n답변 생성"),
    ]
    x = 0.45
    for i, (num, title, body) in enumerate(pipeline):
        round_rect(s, x, 1.55, 2.2, 2.7, SOFT)
        oval = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x + 0.85), Inches(1.75), Inches(0.5), Inches(0.5))
        fill(oval, DARK)
        oval.shadow.inherit = False
        tf = oval.text_frame
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        try:
            oval.text_frame._txBody.bodyPr.set("anchor", "ctr")
        except Exception:
            pass
        r = tf.paragraphs[0].add_run()
        r.text = num
        set_run(r, 14, bold=True, color=GOLD)
        add_text(s, x + 0.1, 2.4, 2.0, 0.4, title, 16, bold=True, color=BLACK, align=PP_ALIGN.CENTER)
        add_text(s, x + 0.12, 2.85, 1.96, 1.1, body, 12, color=MUTED, align=PP_ALIGN.CENTER)
        if i < len(pipeline) - 1:
            add_text(s, x + 2.08, 2.55, 0.4, 0.4, "→", 20, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
        x += 2.55

    add_lines(
        s, 0.55, 4.55, 12.2, 2.2,
        [
            ("한 줄로 보면", True, BLACK, 15),
            ("database/mcm_products_v2.json  →  rag/documents.jsonl  →  vector_db  →  질문 검색  →  OpenAI / Qwen 답변", False, INK, 16),
            ("검색은 제품 modelCode로 후보를 좁힌 뒤, 질문과 문서 벡터의 유사도로 관련 문서를 고릅니다.", False, MUTED, 14),
        ],
        spacing=10,
    )

    # --------------------------------------------------
    # 4. Database
    # --------------------------------------------------
    s = new_slide(prs)
    header(s, "01  DATABASE", "제품 Database")
    footer(s, 4)

    labeled_box(
        s, 0.55, 1.4, 6.1, 2.55,
        "database/mcm_products_v2.json",
        "MCM 한국 공식 제품 페이지 기준 시드 데이터입니다.\n가격·재고 같은 변동 정보는 넣지 않고, 제품 설명·소재·사양·관리법을 중심으로 구성합니다.",
    )
    labeled_box(
        s, 6.9, 1.4, 5.9, 2.55,
        "현재 규모",
        "제품 20개. 가방, 지갑, 의류, 신발 등 카테고리별 공식 정보를 담고 있습니다.\n예: Tracy 비세토스 크로스바디 (MWRAAME01CO001)",
    )

    fields = [
        ("기본 정보", "제품명, modelCode, 컬러, 사이즈"),
        ("설명", "개요, 헤리티지, 핵심 키워드"),
        ("사양", "소재, 크기, 하드웨어, 기능"),
        ("관리", "보관, 세탁, 주의사항"),
    ]
    x = 0.55
    for title, body in fields:
        round_rect(s, x, 4.2, 2.95, 2.4, SOFT)
        add_text(s, x + 0.18, 4.4, 2.6, 0.4, title, 16, bold=True, color=BLACK)
        add_text(s, x + 0.18, 4.9, 2.6, 1.3, body, 13, color=MUTED)
        x += 3.15

    # --------------------------------------------------
    # 5. RAG + Vector
    # --------------------------------------------------
    s = new_slide(prs)
    header(s, "02–03  RAG  ·  VECTOR", "문서로 나눈 뒤, 벡터로 만듭니다")
    footer(s, 5)

    round_rect(s, 0.55, 1.4, 6.0, 5.4, SOFT)
    add_text(s, 0.8, 1.6, 5.5, 0.35, "RAG 문서 생성", 18, bold=True, color=BLACK)
    add_text(s, 0.8, 2.05, 5.5, 0.28, "build_rag_documents.py", 12, bold=True, color=GOLD)
    add_lines(
        s, 0.8, 2.5, 5.5, 4.0,
        [
            "제품 JSON을 통째로 쓰지 않고, 주제별 짧은 문서로 분리합니다.",
            "",
            "overview  ·  제품 개요",
            "description_heritage  ·  디자인 / 헤리티지",
            "specifications  ·  사양",
            "materials  ·  소재",
            "features  ·  기능",
            "care  ·  보관 / 세탁 / 주의사항",
            "",
            "결과물  →  rag/documents.jsonl",
        ],
        size=14,
        spacing=4,
    )

    round_rect(s, 6.8, 1.4, 6.0, 5.4, DARK)
    add_text(s, 7.05, 1.6, 5.5, 0.35, "벡터화", 18, bold=True, color=WHITE)
    add_text(s, 7.05, 2.05, 5.5, 0.28, "build_vector_db.py", 12, bold=True, color=GOLD)
    add_lines(
        s, 7.05, 2.5, 5.5, 4.0,
        [
            ("문장 의미를 숫자 벡터로 바꿉니다.", False, WHITE, 14),
            ("", False, WHITE, 10),
            ("Embedding 모델", True, GOLD, 13),
            ("BAAI/bge-m3", False, WHITE, 14),
            ("", False, WHITE, 10),
            ("Vector Index", True, GOLD, 13),
            ("FAISS  ·  IndexFlatIP", False, WHITE, 14),
            ("정규화 후 내적 = 코사인 유사도", False, RGBColor(0xC8, 0xC3, 0xBA), 13),
            ("", False, WHITE, 10),
            ("결과물  →  vector_db/", False, WHITE, 14),
            ("documents.json  ·  index.faiss", False, RGBColor(0xC8, 0xC3, 0xBA), 13),
        ],
        spacing=4,
    )

    # --------------------------------------------------
    # 6. Search + Agent
    # --------------------------------------------------
    s = new_slide(prs)
    header(s, "04–05  SEARCH  ·  AGENT", "관련 문서를 찾고, 그 내용으로 답합니다")
    footer(s, 6)

    search_steps = [
        ("1", "질문 입력", "고객 질문 + 제품 modelCode"),
        ("2", "후보 축소", "해당 제품 문서만 남김"),
        ("3", "벡터 비교", "질문 벡터와 문서 벡터의 유사도 계산"),
        ("4", "문서 선택", "관련도 높은 상위 3개(Top-K) 사용"),
        ("5", "Agent 답변", "선택한 문서를 근거로 답변 생성"),
    ]
    y = 1.4
    for num, title, body in search_steps:
        oval = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.55), Inches(y), Inches(0.42), Inches(0.42))
        fill(oval, DARK)
        oval.shadow.inherit = False
        tf = oval.text_frame
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        try:
            oval.text_frame._txBody.bodyPr.set("anchor", "ctr")
        except Exception:
            pass
        r = tf.paragraphs[0].add_run()
        r.text = num
        set_run(r, 13, bold=True, color=GOLD)
        add_text(s, 1.15, y - 0.02, 5.4, 0.22, title, 16, bold=True, color=BLACK)
        add_text(s, 1.15, y + 0.28, 5.4, 0.28, body, 13, color=MUTED)
        y += 1.05

    round_rect(s, 7.15, 1.4, 5.6, 5.4, SOFT)
    add_text(s, 7.4, 1.65, 5.15, 0.4, "AI Agent 원칙", 18, bold=True, color=BLACK)
    add_lines(
        s, 7.4, 2.2, 5.15, 4.2,
        [
            "역할",
            "MCM 프리미엄 브랜드 컨시어지",
            "",
            "근거",
            "검색된 Product Information만 사용",
            "",
            "제한",
            "확인되지 않은 소재·관리법·사양은 만들지 않음",
            "",
            "부족할 때",
            "공식 제품 페이지 또는 고객 서비스로 안내",
        ],
        size=14,
        spacing=3,
    )

    # --------------------------------------------------
    # 7. OpenAI / Qwen
    # --------------------------------------------------
    s = new_slide(prs)
    header(s, "LLM", "지금은 OpenAI, 이후에는 Qwen")
    footer(s, 7)

    round_rect(s, 0.55, 1.4, 6.0, 5.4, DARK)
    add_text(s, 0.8, 1.65, 5.5, 0.28, "현재", 12, bold=True, color=GOLD)
    add_text(s, 0.8, 2.0, 5.5, 0.45, "OpenAI", 32, bold=True, color=WHITE)
    add_lines(
        s, 0.8, 2.65, 5.5, 3.8,
        [
            ("chat_openai.py  ·  chat_ui_openai.py  ·  api.py", False, RGBColor(0xC8, 0xC3, 0xBA), 13),
            ("", False, WHITE, 8),
            ("기본 모델  gpt-4o-mini", False, WHITE, 16),
            ("설정  config/openai_config.json", False, WHITE, 16),
            ("API  POST /chat/stream", False, WHITE, 16),
            ("", False, WHITE, 8),
            ("검색 파이프라인은 그대로 두고,", False, RGBColor(0xD6, 0xD0, 0xC6), 14),
            ("답변 생성만 OpenAI가 담당합니다.", False, RGBColor(0xD6, 0xD0, 0xC6), 14),
        ],
        spacing=6,
    )

    round_rect(s, 6.8, 1.4, 6.0, 5.4, SOFT)
    add_text(s, 7.05, 1.65, 5.5, 0.28, "확장", 12, bold=True, color=GOLD)
    add_text(s, 7.05, 2.0, 5.5, 0.45, "Qwen", 32, bold=True, color=BLACK)
    add_lines(
        s, 7.05, 2.65, 5.5, 3.8,
        [
            ("qwen/chat.py  ·  qwen/chat_ui.py", False, MUTED, 13),
            ("", False, MUTED, 8),
            ("모델  Qwen/Qwen3.5-2B", False, BLACK, 16),
            ("방식  로컬 LLM  ·  4bit 양자화", False, BLACK, 16),
            ("검색  동일한 RAG / 벡터 검색", False, BLACK, 16),
            ("", False, MUTED, 8),
            ("Database → RAG → Vector → Search는", False, MUTED, 14),
            ("공유하고, LLM만 교체하면 됩니다.", False, MUTED, 14),
        ],
        spacing=6,
    )

    # --------------------------------------------------
    # 8. Summary
    # --------------------------------------------------
    s = new_slide(prs)
    header(s, "SUMMARY", "한 장으로 정리")
    footer(s, 8)

    rows = [
        ("Database", "database/", "공식 제품 JSON을 원본으로 사용"),
        ("RAG", "rag/", "제품을 주제별 문서로 분할"),
        ("Vector", "vector_db/", "bge-m3로 embedding 후 FAISS에 저장"),
        ("Search", "chat / test_search", "질문 벡터와 비교해 관련 문서 Top-3 선택"),
        ("Agent", "OpenAI → Qwen", "검색된 문서를 근거로 컨시어지가 답변"),
    ]
    y = 1.4
    for i, (title, path, desc) in enumerate(rows):
        bg = SOFT if i % 2 == 0 else WHITE
        round_rect(s, 0.55, y, 12.2, 0.95, bg)
        add_text(s, 0.8, y + 0.22, 2.4, 0.5, title, 18, bold=True, color=BLACK, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, 3.3, y + 0.22, 3.3, 0.5, path, 14, color=GOLD, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, 6.7, y + 0.22, 5.7, 0.5, desc, 15, color=INK, anchor=MSO_ANCHOR.MIDDLE)
        y += 1.05

    out = "MCM_AI_Concierge.pptx"
    prs.save(out)
    print(out)


if __name__ == "__main__":
    main()
